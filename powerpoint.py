import scaper
import os
from pptx import Presentation
from pptx.util import Inches, Pt

FILE_EXTENTSION = ".txt"
FONT_SIZE = Pt(25)


def loadSongs(directory):
    files = []
    if directory != "":
        directory += os.path.sep
    directory = scaper.OUTPUT_PATH + directory
    if os.path.exists(directory):
        for file in os.listdir(directory):
            if file.endswith(FILE_EXTENTSION):
                if not file.startswith("errors"):
                    files.append(directory + file)

    else:
        print("error: invalid directory: " + directory)

    return files


def openSong(song):
    lines = []
    song_list = []
    credits_stuff = []
    with open(song) as f:
        lines = f.readlines()

    if len(lines) != 0:
        count = 0
        if lines[0].startswith("Intro"):
            count = 1
        currentStart = 0
        newName = ""
        startFound = False
        while count < len(lines)-1:
            if lines[count] == "\n":
                if startFound:
                    song_list.append([newName, lines[currentStart: count]])
                    newName = lines[count+1]
                    count += 2
                    currentStart = count
            elif not startFound:
                startFound = True
                newName = lines[count]
                currentStart = count + 1

            count += 1

            if lines[len(lines)-2].startswith("Author"):
                credits_stuff.append(lines[len(lines)-2].split(":")[-1])
            if lines[len(lines)-1].startswith("CCLI"):
                credits_stuff.append(lines[len(lines)-1].split(":")[-1])
    else:
        print("error: " + song)

    return song, song_list, credits_stuff


def createPowerpoint(title, song_list, details):
    prs = Presentation()
    left = top = width = height = Inches(1)
    blank_slide_layout = prs.slide_layouts[6]
    for song in song_list:
        s_details = song[0] + "Writers: " + details[0] + "CCLI: " + details[1]
        slide = prs.slides.add_slide(blank_slide_layout)

        p_lines = slide.shapes.add_textbox(left, top, Inches(8), height).text_frame.add_paragraph()
        p_lines.text = "".join(song[1])
        p_lines.font.size = FONT_SIZE

        p_details = slide.shapes.add_textbox(Inches(5), Inches(6.5), width, height).text_frame.add_paragraph()
        p_details.text = s_details
        p_details.font.size = Pt(12)

        slide.notes_slide.notes_text_frame.text = s_details

    prs.save(title.replace(FILE_EXTENTSION, "") + ".pptx")


def runDir(directory):
    songs = loadSongs(directory)
    for song in songs:
        title, lines, details = openSong(song)
        createPowerpoint(title, lines, details)


def main():
    import time
    s = time.time()

    for file in ["hymns", "notWellKnown", "modernWorship"]:
        runDir(file)

    print("finished: " + str((time.time() - s)))
if __name__ == '__main__':
    main()